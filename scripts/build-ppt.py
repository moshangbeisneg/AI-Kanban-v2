#!/usr/bin/env python3
"""Generate AI Kanban business plan PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x25, 0x60, 0xEB)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
DARK = RGBColor(0x1A, 0x1D, 0x23)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF5, 0xF7)
LIGHT_BLUE = RGBColor(0xEE, 0xF2, 0xFF)
LIGHT_GRAY = RGBColor(0xE2, 0xE4, 0xE8)
GREEN = RGBColor(0x10, 0xB9, 0x81)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_bullet_block(slide, items, left, top, width, height, size=13, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "\u2022  " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return txBox

# ===== Slide 1: Cover =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, 0, 3.3, 10, 0.04, fill_color=GREEN)
add_text(slide, "AI Kanban", 1, 1.5, 8, 1.2, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "\u667a\u80fd\u4efb\u52a1\u7ba1\u7406\u5e73\u53f0 \u00b7 \u5546\u4e1a\u8ba1\u5212\u4e66", 1, 2.8, 8, 0.6, size=20, color=RGBColor(0xBE, 0xD7, 0xFE), align=PP_ALIGN.CENTER)
add_text(slide, "\u8ba9\u56e2\u961f\u534f\u4f5c\u66f4\u9ad8\u6548\uff0c\u8ba9\u9879\u76ee\u7ba1\u7406\u66f4\u667a\u80fd", 1, 3.6, 8, 0.5, size=15, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, "\u521b\u4e1a\u56e2\u961f \u00b7 2026\u5e747\u6708", 1, 5.5, 8, 0.5, size=13, color=RGBColor(0x60, 0xA5, 0xFA), align=PP_ALIGN.CENTER)

# ===== Slide 2: Team =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u5173\u4e8e\u6211\u4eec", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
add_text(slide, "\u521d\u521b\u56e2\u961f \u00b7 \u6280\u672f\u521b\u65b0\u9a71\u52a8 \u00b7 \u6df1\u8015\u56e2\u961f\u534f\u4f5c\u9886\u57df", 0.6, 1.4, 8.8, 0.5, size=13, color=GRAY)

# motto
add_shape(slide, 0.6, 2.0, 8.8, 0.7, fill_color=LIGHT_BLUE)
add_text(slide, '\u201c\u6211\u4eec\u76f8\u4fe1\uff0c\u6700\u597d\u7684\u5de5\u5177\u5e94\u8be5\u8ba9\u56e2\u961f\u4e13\u6ce8\u4e8e\u521b\u9020\uff0c\u800c\u4e0d\u662f\u7ba1\u7406\u6d41\u7a0b\u3002\u201d', 0.8, 2.15, 8.4, 0.4, size=12, color=DARK_BLUE, align=PP_ALIGN.CENTER)

members = [
    ("\u5f20\u660e", "CEO / \u4ea7\u54c1", "\u8fde\u7eed\u521b\u4e1a\u8005\uff0c\u524d\u67d0SaaS\u516c\u53f8\u4ea7\u54c1\u7ecf\u7406\n5\u5e74\u9879\u76ee\u7ba1\u7406\u7ecf\u9a8c\uff0c\u4e3b\u5bfc\u8fc73\u6b3e\u4ea7\u54c1\u4ece0\u52301"),
    ("\u674e\u5a77", "CTO / \u6280\u672f", "\u5168\u6808\u5de5\u7a0b\u5e08\uff0c\u64c5\u957f\u524d\u7aef\u67b6\u6784\n\u66fe\u4efb\u67d0\u4e92\u8054\u7f51\u516c\u53f8\u6280\u672fLead"),
    ("\u738b\u6d69", "COO / \u8fd0\u8425", "\u5e02\u573a\u8fd0\u8425\u80cc\u666f\uff0c\u64c5\u957f\u7528\u6237\u589e\u957f\n\u4e4b\u524d\u8d1f\u8d23\u67d0\u5de5\u5177\u7c7b\u4ea7\u54c1\u83b7\u5ba2"),
]
for i, (name, role, bio) in enumerate(members):
    x = 0.5 + i * 3.15
    add_shape(slide, x, 3.0, 2.95, 3.2, fill_color=LIGHT_BG)
    add_text(slide, name, x + 0.2, 3.2, 2.55, 0.35, size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, role, x + 0.2, 3.55, 2.55, 0.3, size=11, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, bio, x + 0.2, 4.0, 2.55, 1.8, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)

add_shape(slide, 0.6, 6.5, 8.8, 0.7, fill_color=LIGHT_BLUE)
add_text(slide, "\u56e2\u961f\u534f\u4f5c\u65b9\u5f0f\uff1a\u8fdc\u7a0b\u5206\u5e03\u5f0f\u529e\u516c \u00b7 \u6bcf\u65e5\u7ad9\u4f1a\u540c\u6b65 \u00b7 \u53cc\u5468\u8fed\u4ee3\u4ea4\u4ed8 \u00b7 \u5168\u5458\u53c2\u4e0e\u4ea7\u54c1\u51b3\u7b56", 0.8, 6.62, 8.4, 0.35, size=11, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ===== Slide 3: Pain Points & Demand =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u5e02\u573a\u75db\u70b9\u4e0e\u9700\u6c42\u5224\u65ad", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
pains = [
    ("\u26d4", "\u4fe1\u606f\u5b64\u5c9b\u4e25\u91cd", "\u4efb\u52a1\u5206\u6563\u5728\u5fae\u4fe1\u3001\u90ae\u4ef6\u3001Excel\u4e2d\uff0c\u65e0\u6cd5\u7edf\u4e00\u8ffd\u8e2a\u8fdb\u5c55"),
    ("\u23f0", "\u4f18\u5148\u7ea7\u7ba1\u7406\u6df7\u4e71", "\u7d27\u6025\u4e0d\u91cd\u8981\u7684\u4e8b\u52a1\u6324\u5360\u6838\u5fc3\u5de5\u4f5c\uff0c\u56e2\u961f\u75b2\u4e8e\u6551\u706b"),
    ("\U0001f4ca", "\u51b3\u7b56\u7f3a\u4e4f\u6570\u636e\u652f\u6491", "\u9879\u76ee\u8fdb\u5c55\u5168\u51ed\u611f\u89c9\uff0c\u98ce\u9669\u548c\u74f6\u9888\u9760\u4eba\u5de5\u53d1\u73b0"),
    ("\U0001f4b0", "\u73b0\u6709\u5de5\u5177\u6210\u672c\u9ad8", "Jira/Asana\u7b49\u4e13\u4e1a\u5de5\u5177\u5bf9\u5c0f\u578b\u56e2\u961f\u8d1f\u62c5\u8fc7\u91cd"),
]
for i, (emoji, title, desc) in enumerate(pains):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 4.7
    y = 1.6 + row * 2.2
    add_shape(slide, x, y, 4.5, 1.9, fill_color=LIGHT_BG)
    add_text(slide, emoji, x + 0.2, y + 0.15, 0.5, 0.4, size=22)
    add_text(slide, title, x + 0.8, y + 0.15, 3.5, 0.35, size=14, bold=True, color=DARK_BLUE)
    add_text(slide, desc, x + 0.2, y + 0.65, 4.1, 1.0, size=11, color=GRAY)
add_shape(slide, 0.5, 6.1, 9.0, 1.1, fill_color=LIGHT_BLUE, line_color=BLUE, line_width=1.5)
add_text(slide, "\u6211\u4eec\u7684\u5224\u65ad\uff1a", 0.7, 6.2, 8.6, 0.35, size=13, bold=True, color=BLUE)
add_text(slide, "\u4e2d\u56fd\u6709\u8d854000\u4e07\u4e2d\u5c0f\u5fae\u4f01\u4e1a\uff0c\u5176\u4e2d\u7edd\u5927\u591a\u6570\u4ecd\u5728\u4f7f\u7528\u788e\u7247\u5316\u5de5\u5177\u7ba1\u7406\u56e2\u961f\u4efb\u52a1\u3002\u4e00\u4e2a\u201c\u8f7b\u91cf + AI + \u4eb2\u6c11\u4ef7\u683c\u201d\u7684\u667a\u80fd\u770b\u677f\u4ea7\u54c1\u6709\u660e\u786e\u7684\u5207\u5165\u673a\u4f1a\u3002", 0.7, 6.55, 8.6, 0.55, size=11, color=DARK)

# ===== Slide 4: Product Solution =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u4ea7\u54c1\u89e3\u51b3\u65b9\u6848", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
add_text(slide, "AI Kanban \u2014\u2014 \u4e00\u4e2a\u81ea\u5e26 AI \u52a9\u624b\u7684\u667a\u80fd\u4efb\u52a1\u770b\u677f", 0.6, 1.4, 8.8, 0.4, size=14, color=GRAY)
features = [
    ("\U0001f4cb \u53ef\u89c6\u5316\u770b\u677f", "\u62d6\u62fd\u5f0f\u4e09\u5217\u5e03\u5c40\uff08\u5f85\u529e/\u8fdb\u884c\u4e2d/\u5df2\u5b8c\u6210\uff09\uff0c\u4efb\u52a1\u6d41\u8f6c\u4e00\u76ee\u4e86\u7136"),
    ("\U0001f916 AI \u667a\u80fd\u5206\u6790", "\u5185\u7f6e\u89c4\u5219\u5f15\u64ce\uff0c\u81ea\u52a8\u68c0\u6d4b\u8d85\u671f\u4efb\u52a1\u3001\u9ad8\u4f18\u5806\u79ef\u3001\u5b8c\u6210\u7387\uff0c\u53ef\u9009\u914d\u7f6e DeepSeek API"),
    ("\U0001f4ac AI \u5bf9\u8bdd\u52a9\u624b", "\u96c6\u6210 DeepSeek API\uff0c\u53ef\u76f4\u63a5\u4e0eAI\u8ba8\u8bba\u4efb\u52a1\u3001\u83b7\u53d6\u4f18\u5148\u7ea7\u5efa\u8bae\u548c\u6392\u671f\u63a8\u8350"),
    ("\U0001f50d \u591a\u7ef4\u7b5b\u9009", "\u6309\u4f18\u5148\u7ea7\u3001\u5173\u952e\u8bcd\u3001\u6807\u7b7e\u7ec4\u5408\u7b5b\u9009\uff0c\u6279\u91cf\u64cd\u4f5c\u652f\u6301\u9ad8\u6548\u7ba1\u7406"),
    ("\U0001f4c8 \u6570\u636e\u7edf\u8ba1", "\u5b8c\u6210\u7387\u3001\u4f18\u5148\u7ea7\u5206\u5e03\u3001\u5404\u5217\u4efb\u52a1\u6570\u53ef\u89c6\u5316\uff0c\u652f\u6301\u5bfc\u51fa JSON/CSV/PDF"),
    ("\u23f1 \u65f6\u95f4\u7ebf\u89c6\u56fe", "\u6309\u622a\u6b62\u65e5\u671f\u6392\u5e8f\u7684\u4efb\u52a1\u65f6\u95f4\u7ebf\uff0c\u6e05\u6670\u628a\u63e1\u9879\u76ee\u8282\u594f"),
]
for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 4.7
    y = 1.8 + row * 1.5
    add_shape(slide, x, y, 4.5, 1.25, fill_color=LIGHT_BG if i % 2 == 0 else WHITE, line_color=LIGHT_GRAY, line_width=0.5)
    add_text(slide, title, x + 0.2, y + 0.08, 4.1, 0.35, size=13, bold=True, color=DARK_BLUE)
    add_text(slide, desc, x + 0.2, y + 0.5, 4.1, 0.6, size=10.5, color=GRAY)
add_shape(slide, 0.5, 6.0, 9.0, 1.2, fill_color=LIGHT_BLUE)
add_text(slide, "\u6280\u672f\u4eae\u70b9\uff1a\u7eaf\u524d\u7aef\u5e94\u7528\uff08HTML+CSS+JS\uff09\uff0clocalStorage\u6301\u4e45\u5316\uff0c\u652f\u6301\u79bb\u7ebf\u4f7f\u7528 \u00b7 \u53ef\u76f4\u63a5\u6253\u5f00index.html\u8fd0\u884c\uff0c\u96f6\u90e8\u7f72\u6210\u672c \u00b7 AI\u6a21\u5757\u53ef\u79bb\u7ebf\u8fd0\u884c\u89c4\u5219\u5f15\u64ce\uff0c\u914d\u7f6eAPI Key\u540e\u63a5\u5165DeepSeek", 0.7, 6.1, 8.6, 1.0, size=10, color=DARK_BLUE)

# ===== Slide 5: Usability & Advancement =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u4ea7\u54c1\u53ef\u7528\u6027\u4e0e\u5148\u8fdb\u6027", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
add_text(slide, "\u53ef\u7528\u6027\u8bbe\u8ba1", 0.6, 1.35, 4.3, 0.4, size=16, bold=True, color=GREEN)
usability = [
    "\u6253\u5f00\u5373\u7528\uff1a\u65e0\u9700\u6ce8\u518c\u3001\u65e0\u9700\u5b89\u88c5\u3001\u6d4f\u89c8\u5668\u6253\u5f00index.html\u5373\u53ef",
    "\u76f4\u89c9\u64cd\u4f5c\uff1a\u62d6\u62fd\u79fb\u52a8\u4efb\u52a1\uff0c\u5361\u7247\u4fe1\u606f\u4e00\u76ee\u4e86\u7136",
    "\u5373\u65f6\u53cd\u9988\uff1a\u6240\u6709\u64cd\u4f5c\u5b9e\u65f6\u751f\u6548\uff0c\u9884\u89c8\u7ed3\u679c\u65e0\u9700\u5237\u65b0",
    "\u6570\u636e\u5b89\u5168\uff1a\u6240\u6709\u6570\u636e\u5b58\u50a8\u5728\u6d4f\u89c8\u5668\u672c\u5730\uff0c\u4e0d\u4e0a\u4f20\u670d\u52a1\u5668",
    "\u591a\u7aef\u9002\u914d\uff1a\u5728PC\u548c\u5e73\u677f\u6d4f\u89c8\u5668\u4e0a\u5747\u80fd\u6b63\u5e38\u4f7f\u7528",
    "\u5b66\u4e60\u6210\u672c\u6781\u4f4e\uff1a\u6838\u5fc3\u64cd\u4f5c\u4ec5\u9700\u62d6\u62fd\u548c\u70b9\u51fb\uff0c\u65b0\u624b5\u5206\u949f\u4e0a\u624b",
]
add_bullet_block(slide, usability, 0.6, 1.85, 4.5, 3.5, size=11, color=DARK)
add_text(slide, "\u6280\u672f\u5148\u8fdb\u6027", 5.2, 1.35, 4.3, 0.4, size=16, bold=True, color=BLUE)
advance = [
    "\u5185\u5d4c AI \u5206\u6790\u5f15\u64ce\uff1a\u65e0\u9700\u4f9d\u8d56\u4e91\u7aef\uff0c\u672c\u5730\u5373\u53ef\u63d0\u4f9b\u667a\u80fd\u6d1e\u5bdf",
    "\u652f\u6301 DeepSeek API\uff1a\u53ef\u9009\u589e\u5f3aAI\u80fd\u529b\uff0c\u652f\u6301\u5bf9\u8bdd\u5f0f\u4efb\u52a1\u7ba1\u7406",
    "\u539f\u751f\u62d6\u62fd API\uff1a\u65e0\u9700\u7b2c\u4e09\u65b9\u5e93\uff0c\u8f7b\u91cf\u9ad8\u6548",
    "\u4e09\u79cd\u89c6\u56fe\u5207\u6362\uff1a\u770b\u677f/\u65f6\u95f4\u7ebf/\u7edf\u8ba1\uff0c\u6ee1\u8db3\u4e0d\u540c\u7ba1\u7406\u573a\u666f",
    "\u5b8c\u6574\u6570\u636e\u6d41\u6c34\u7ebf\uff1a\u65b0\u5efa\u2192\u7f16\u8f91\u2192\u6d41\u8f6c\u2192\u5b8c\u6210\uff0c\u95ed\u73af\u7ba1\u7406",
    "\u6279\u91cf\u64cd\u4f5c+\u591a\u6761\u4ef6\u7b5b\u9009\uff1a\u9ad8\u6548\u5904\u7406\u5927\u91cf\u4efb\u52a1",
]
add_bullet_block(slide, advance, 5.2, 1.85, 4.5, 3.5, size=11, color=DARK)
add_shape(slide, 0.5, 6.2, 9.0, 1.0, fill_color=LIGHT_BLUE)
add_text(slide, "\u4e0e\u5176\u4ed6\u770b\u677f\u5de5\u5177\u5bf9\u6bd4", 0.7, 6.3, 8.6, 0.3, size=12, bold=True, color=DARK_BLUE)
add_text(slide, "vs Trello\uff1a\u66f4\u8f7b\u91cf\uff0c\u65e0\u9700\u6ce8\u518c \u00b7 vs Jira\uff1a\u96f6\u5b66\u4e60\u6210\u672c\uff0c\u66f4\u9002\u5408\u5c0f\u56e2\u961f \u00b7 vs Notion\uff1a\u4e13\u6ce8\u770b\u677f\u573a\u666f\uff0cAI\u80fd\u529b\u66f4\u805a\u7126 \u00b7 vs Excel\uff1a\u53ef\u89c6\u5316+AI\u8f85\u52a9\uff0c\u6548\u7387\u63d0\u5347\u663e\u8457", 0.7, 6.65, 8.6, 0.45, size=10, color=GRAY)

# ===== Slide 6: Market Prospect =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u5e02\u573a\u524d\u666f", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
mc = [
    ("\U0001f4ca", "\u5168\u7403\u5e02\u573a\u89c4\u6a21", "150\u4ebf\u7f8e\u5143", "\u9879\u76ee\u7ba1\u7406\u8f6f\u4ef6\u5e02\u573a\uff082026\u5e74\uff09", "\u5e74\u5747\u589e\u957f\u738712%\uff0cAI+\u9879\u76ee\u7ba1\u7406\u7ec6\u5206\u589e\u901f25%"),
    ("\U0001f3af", "\u76ee\u6807\u7528\u6237\u753b\u50cf", "4000\u4e07+\u4e2d\u5c0f\u5fae\u4f01\u4e1a", "\u6838\u5fc3\uff1a10-50\u4eba\u521b\u4e1a/\u4e2d\u5c0f\u56e2\u961f", "\u8feb\u5207\u9700\u8981\u8f7b\u91cf\u3001\u667a\u80fd\u3001\u4f4e\u6210\u672c\u7684\u7ba1\u7406\u5de5\u5177"),
    ("\U0001f4c8", "\u589e\u957f\u9a71\u52a8\u529b", "\u4e09\u5927\u8d8b\u52bf\u53e0\u52a0", "\u8fdc\u7a0b\u529e\u516c\u666e\u53ca + AI\u5de5\u5177\u5e73\u6c11\u5316 + \u4e2d\u5c0f\u4f01\u4e1a\u6570\u5b57\u5316\u8f6c\u578b", "\u5e02\u573a\u7a97\u53e3\u671f\u6b63\u5728\u5c55\u5f00"),
]
for i, (emoji, title, value, sub, desc) in enumerate(mc):
    x = 0.4 + i * 3.2
    add_shape(slide, x, 1.5, 3.0, 3.8, fill_color=LIGHT_BG)
    add_text(slide, emoji, x + 0.2, 1.65, 0.4, 0.4, size=24)
    add_text(slide, title, x + 0.2, 2.1, 2.6, 0.35, size=13, bold=True, color=DARK_BLUE)
    add_text(slide, value, x + 0.2, 2.5, 2.6, 0.4, size=18, bold=True, color=BLUE)
    add_text(slide, sub, x + 0.2, 3.0, 2.6, 0.5, size=10.5, color=GRAY)
    add_text(slide, desc, x + 0.2, 3.7, 2.6, 1.3, size=10.5, color=DARK)
add_shape(slide, 0.5, 5.6, 9.0, 1.6, fill_color=LIGHT_BLUE)
add_text(slide, "\u6211\u4eec\u7684\u7ade\u4e89\u4f18\u52bf", 0.7, 5.7, 8.6, 0.3, size=13, bold=True, color=BLUE)
add_text(slide, "\u2022 \u5148\u53d1\u4f18\u52bf\uff1a\u56fd\u5185\u56e2\u961f\u6025\u9700\u201cAI+\u770b\u677f\u201d\u7684\u7ed3\u5408\u4ea7\u54c1\uff0c\u5e02\u573a\u5c1a\u5728\u65e9\u671f\n\u2022 \u6210\u672c\u4f18\u52bf\uff1a\u7eaf\u524d\u7aef\u67b6\u6784\uff0c\u670d\u52a1\u5668\u6210\u672c\u6781\u4f4e\uff0c\u53ef\u652f\u6491\u514d\u8d39\u7b56\u7565\u83b7\u5ba2\n\u2022 \u8fed\u4ee3\u4f18\u52bf\uff1a\u8f7b\u91cf\u67b6\u6784\u652f\u6301\u5feb\u901f\u8fed\u4ee3\uff0c\u53ef\u6839\u636e\u7528\u6237\u53cd\u9988\u5468\u7ea7\u66f4\u65b0", 0.7, 6.05, 8.6, 1.0, size=11, color=DARK)

# ===== Slide 7: User Reach & Conversion =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u7528\u6237\u89e6\u8fbe\u4e0e\u8f6c\u5316\u7b56\u7565", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
add_text(slide, "\u83b7\u5ba2\u6e20\u9053", 0.6, 1.35, 4.3, 0.4, size=16, bold=True, color=GREEN)
channels = [
    "\u5f00\u6e90\u793e\u533a\uff08GitHub\uff09\uff1a\u514d\u8d39\u5f00\u6e90\uff0c\u5438\u5f15\u5f00\u53d1\u8005\u8d21\u732e\u4e0e\u4f20\u64ad\uff0c\u7b2c\u4e00\u6279\u79cd\u5b50\u7528\u6237",
    "\u79d1\u6280\u793e\u533a\uff1a\u53d1\u5e03\u5230V2EX\u3001\u5373\u523b\u3001\u5c11\u6570\u6d3e\u7b49\u79d1\u6280\u793e\u533a",
    "\u53e3\u7891\u88c2\u53d8\uff1a\u5185\u5efa\u5bfc\u51fa\u5206\u4eab\u529f\u80fd\uff0c\u7528\u6237\u53ef\u4e00\u952e\u5206\u4eab\u770b\u677f\u5feb\u7167",
    "\u5185\u5bb9\u8425\u9500\uff1a\u8f93\u51fa\u201cAI+\u9879\u76ee\u7ba1\u7406\u201d\u76f8\u5173\u6559\u7a0b\u3001\u6700\u4f73\u5b9e\u8df5\u6587\u7ae0",
    "\u6821\u56ed\u6e20\u9053\uff1a\u9762\u5411\u8ba1\u7b97\u673a\u76f8\u5173\u8bfe\u7a0b\u7684\u5b66\u751f\u63a8\u5e7f\uff0c\u57f9\u517b\u65e9\u671f\u4f7f\u7528\u4e60\u60ef",
]
add_bullet_block(slide, channels, 0.6, 1.85, 4.5, 3.5, size=11, color=DARK)
add_text(slide, "\u8f6c\u5316\u8def\u5f84", 5.2, 1.35, 4.3, 0.4, size=16, bold=True, color=BLUE)
conversion = [
    "\u5f00\u6e90\u514d\u8d39\u7248 \u2192 \u4f53\u9a8c\u6838\u5fc3\u529f\u80fd \u2192 \u4ea7\u751f\u4f9d\u8d56 \u2192 \u5347\u7ea7\u4e13\u4e1a\u7248",
    "Freemium\u6a21\u578b\uff1a\u57fa\u7840\u529f\u80fd\u514d\u8d39\uff0c\u9ad8\u7ea7AI\u548c\u5206\u6790\u529f\u80fd\u8ba2\u9605",
    "\u5173\u952e\u8f6c\u5316\u8282\u70b9\uff1a\u521b\u5efa\u7b2c10\u4e2a\u4efb\u52a1\u65f6\u63d0\u793a\u5347\u7ea7",
    "\u56e2\u961f\u9080\u8bf7\u673a\u5236\uff1a1\u4eba\u4f7f\u7528\u540e\u9080\u8bf7\u56e2\u961f\u6210\u5458\uff0c\u5f62\u6210\u7f51\u7edc\u6548\u5e94",
    "\u4f01\u4e1aPOC\uff1a\u4e3a\u6709\u9700\u6c42\u7684\u4f01\u4e1a\u63d0\u4f9b15\u5929\u56e2\u961f\u7248\u8bd5\u7528",
]
add_bullet_block(slide, conversion, 5.2, 1.85, 4.5, 3.5, size=11, color=DARK)
add_shape(slide, 0.5, 6.0, 9.0, 1.2, fill_color=LIGHT_BLUE)
add_text(slide, "\u589e\u957f\u6a21\u578b\u9884\u4f30\uff08\u9996\u5e74\uff09", 0.7, 6.1, 8.6, 0.35, size=12, bold=True, color=BLUE)
growth = [
    "\u7b2c1-3\u6708\uff1aGitHub\u5f00\u6e90\u53d1\u5e03\uff0c\u83b7\u53d6 1,000+ Star\uff0c\u79ef\u7d2f 500 \u6d3b\u8dc3\u7528\u6237",
    "\u7b2c4-6\u6708\uff1a\u793e\u533a\u53e3\u7891\u4f20\u64ad\uff0c\u7528\u6237\u589e\u957f\u81f3 5,000\uff0c\u63a8\u51fa\u4e13\u4e1a\u7248\u8ba2\u9605",
    "\u7b2c7-12\u6708\uff1a\u6708\u6d3b\u7528\u6237\u8fbe\u5230 20,000\uff0c\u4e13\u4e1a\u7248\u4ed8\u8d39\u8f6c\u5316\u7387 5%\uff0c\u6708\u5e38\u89c4\u6536\u5165\uff08MRR\uff09\u7ea6 30,000 \u5143",
]
add_bullet_block(slide, growth, 0.7, 6.5, 8.6, 0.65, size=10.5, color=DARK)

# ===== Slide 8: Roadmap & Funding =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, "\u53d1\u5c55\u89c4\u5212\u4e0e\u878d\u8d44\u9700\u6c42", 0.6, 0.4, 8.8, 0.7, size=28, bold=True, color=DARK_BLUE)
add_shape(slide, 0.6, 1.1, 1.2, 0.06, fill_color=BLUE)
phases = [
    ("Phase 1", "\u5f53\u524d - 2026.09", "MVP\u4e0a\u7ebf", "\u6838\u5fc3\u770b\u677f + AI\u89c4\u5219\u5f15\u64ce + \u672c\u5730\u5b58\u50a8\n\u4ea7\u54c1\u6253\u78e8\uff0c\u6536\u96c6\u7b2c\u4e00\u6279\u7528\u6237\u53cd\u9988"),
    ("Phase 2", "2026.10 - 2026.12", "\u56e2\u961f\u7248\u53d1\u5e03", "\u591a\u9879\u76ee\u652f\u6301 + \u5b9e\u65f6\u540c\u6b65 + \u56e2\u961f\u534f\u4f5c\n\u7528\u6237\u589e\u957f\u81f35,000\uff0c\u542f\u52a8\u8ba2\u9605\u6a21\u5f0f"),
    ("Phase 3", "2027.01 - 2027.03", "AI\u589e\u5f3a", "\u6df1\u5ea6AI\u5206\u6790 + \u6570\u636e\u770b\u677f + \u7b2c\u4e09\u65b9\u96c6\u6210API\n\u4ed8\u8d39\u7528\u6237\u8f6c\u5316\uff0c\u76ee\u68075%\u8f6c\u5316\u7387"),
    ("Phase 4", "2027.04 - 2027.06", "\u79fb\u52a8\u7aef+\u5546\u4e1a\u5316", "\u79fb\u52a8\u7aefApp + \u4f01\u4e1a\u7248\u79c1\u6709\u90e8\u7f72\n\u6708\u6d3b20,000\uff0cMRR\u8fbe3\u4e07\u5143"),
]
for i, (phase, period, title, desc) in enumerate(phases):
    y = 1.4 + i * 1.0
    fill = LIGHT_BLUE if i == 0 else LIGHT_BG
    add_shape(slide, 0.6, y, 8.8, 0.85, fill_color=fill, line_color=BLUE if i == 0 else LIGHT_GRAY, line_width=1.5 if i == 0 else 0.5)
    add_text(slide, phase, 0.8, y + 0.05, 1.5, 0.25, size=11, bold=True, color=BLUE)
    add_text(slide, period, 2.3, y + 0.05, 2.5, 0.25, size=9.5, color=GRAY)
    add_text(slide, title, 0.8, y + 0.35, 2.0, 0.25, size=10.5, bold=True, color=DARK)
    add_text(slide, desc, 2.3, y + 0.35, 6.8, 0.45, size=9.5, color=GRAY)
add_shape(slide, 0.5, 5.7, 9.0, 1.6, fill_color=LIGHT_BLUE)
add_text(slide, "\u878d\u8d44\u9700\u6c42", 0.7, 5.8, 8.6, 0.35, size=14, bold=True, color=BLUE)
funding = [
    "\u672c\u8f6e\u878d\u8d44\u76ee\u6807\uff1a50\u4e07\u5143\u4eba\u6c11\u5e01\uff08\u79cd\u5b50\u8f6e\uff09",
    "\u8d44\u91d1\u7528\u9014\uff1a40% \u4ea7\u54c1\u7814\u53d1\uff08\u540e\u7aef+AI\u80fd\u529b\uff09\u300130% \u5e02\u573a\u63a8\u5e7f\u300120% \u56e2\u961f\u5efa\u8bbe\u300110% \u8fd0\u8425\u53ca\u6cd5\u52a1",
    "\u51fa\u8ba9\u80a1\u4efd\uff1a5% - 8%\uff08\u5177\u4f53\u53ef\u534f\u5546\uff09",
    "\u91cc\u7a0b\u7891\u5bf9\u8d4c\uff1a12\u4e2a\u6708\u5185\u5b9e\u73b0 MAU 20,000\uff0cMRR 30,000 \u5143",
]
add_bullet_block(slide, funding, 0.7, 6.15, 8.6, 1.0, size=11, color=DARK)

# ===== Slide 9: Thank You =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, 0, 3.5, 10, 0.04, fill_color=GREEN)
add_text(slide, "\u611f\u8c22\u60a8\u7684\u5173\u6ce8", 1, 1.8, 8, 1.0, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "AI Kanban \u00b7 \u8ba9\u56e2\u961f\u534f\u4f5c\u66f4\u667a\u80fd", 1, 3.0, 8, 0.6, size=16, color=RGBColor(0xBE, 0xD7, 0xFE), align=PP_ALIGN.CENTER)
add_text(slide, "GitHub: https://github.com/moshangbeisneg/AI-Kanban-v2", 1, 4.2, 8, 0.5, size=13, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, "\u671f\u5f85\u4e0e\u60a8\u8fdb\u4e00\u6b65\u4ea4\u6d41\uff01", 1, 5.6, 8, 0.5, size=15, color=RGBColor(0x60, 0xA5, 0xFA), align=PP_ALIGN.CENTER)

# Save
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "AI-Kanban-\u5546\u4e1a\u8ba1\u5212\u4e66.pptx")
prs.save(out_path)
print("PPT generated: " + out_path)
